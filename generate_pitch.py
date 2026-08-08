"""
Generate Aegis Pitch Deck — .pptx (v3 with Slide Transitions & Animations)
Matches AMEX Round 1 Submission Guidelines.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
import os

# ─── Config ──────────────────────────────────────────────────────────────────

OUT_FILE = os.path.join(os.path.dirname(__file__), "Aegis_Pitch_Deck.odp")
IMG_DIR = os.path.join(os.path.dirname(__file__), "docs_images")

IMG_UI   = os.path.join(IMG_DIR, "control_tower_ui.png")
IMG_AUD  = os.path.join(IMG_DIR, "audit_log.png")
IMG_ARCH = os.path.join(IMG_DIR, "architecture_diagram.png")

# Colors
BG_DARK    = RGBColor(0x0B, 0x0F, 0x19)
BG_PANEL   = RGBColor(0x12, 0x17, 0x24)
GOLD       = RGBColor(0xE8, 0xA3, 0x3D)
WHITE      = RGBColor(0xE6, 0xE8, 0xEB)
GRAY       = RGBColor(0x8A, 0x93, 0xA3)
GREEN      = RGBColor(0x4C, 0x9A, 0x6A)
RED        = RGBColor(0xC1, 0x44, 0x3B)
BLUE_ACC   = RGBColor(0x3B, 0x82, 0xF6)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ─── Helpers ─────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_slide_transition(slide, transition_type="fade"):
    """Inject smooth XML slide transition (fade / push / wipe)"""
    if transition_type == "fade":
        xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>'
    elif transition_type == "push":
        xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:push dir="r"/></p:transition>'
    elif transition_type == "wipe":
        xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:wipe dir="r"/></p:transition>'
    else:
        xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>'
    
    slide._element.append(parse_xml(xml))

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_paragraph(text_frame, text, font_size=16, color=WHITE, bold=False,
                  space_before=Pt(6), alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    return p

def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.6), color=GOLD):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color=BG_PANEL):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0x2A, 0x30, 0x3C)
    shape.line.width = Pt(1)
    return shape

def add_slide_number(slide, num, total):
    add_text_box(
        slide, Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.4),
        f"{num} / {total}", font_size=10, color=GRAY, alignment=PP_ALIGN.RIGHT
    )

TOTAL_SLIDES = 10

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_transition(slide, "fade")

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = GOLD
shape.line.fill.background()

add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
             "AEGIS", font_size=60, color=GOLD, bold=True)

add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8),
             "Governance Control for Financial AI Agents", font_size=32, color=WHITE, bold=True)

add_text_box(slide, Inches(1.5), Inches(4.3), Inches(10), Inches(0.5),
             "Real-time policy enforcement  •  Cryptographic audit trail  •  Fleet-wide kill switches",
             font_size=16, color=GRAY)

add_text_box(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.4),
             "American Express Hackathon 2026  |  Governance Layer for Financial Agents",
             font_size=14, color=GRAY)

add_slide_number(slide, 1, TOTAL_SLIDES)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_transition(slide, "fade")

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "THE PROBLEM", font_size=14, color=GOLD, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(1.0),
             "Financial AI agents are being deployed without guardrails.",
             font_size=36, color=WHITE, bold=True)

problems = [
    ("⚠  Runaway Spend", "Autonomous agents can exceed daily budgets when actions execute before checks. Institutions face unbounded financial liability."),
    ("🔓  No Policy Enforcement", "Agents act first and get audited later. No mandatory interceptor blocks non-compliant actions prior to execution."),
    ("📋  Audit Trail Gaps", "Post-hoc logging is alterable and unverifiable. Regulators cannot prove historical record integrity without cryptographic proof."),
]

for i, (title, desc) in enumerate(problems):
    left = Inches(0.8 + i * 4.0)
    card = add_rounded_rect(slide, left, Inches(2.8), Inches(3.6), Inches(3.8))

    add_text_box(slide, left + Inches(0.3), Inches(3.0), Inches(3.0), Inches(0.6),
                 title, font_size=20, color=GOLD, bold=True)

    add_text_box(slide, left + Inches(0.3), Inches(3.8), Inches(3.0), Inches(2.4),
                 desc, font_size=15, color=GRAY)

add_slide_number(slide, 2, TOTAL_SLIDES)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — OUR SOLUTION (with real Control Tower UI screenshot)
# ═════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_transition(slide, "fade")

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.5),
             "OUR SOLUTION", font_size=14, color=GOLD, bold=True)

add_text_box(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.6),
             "Aegis — Control Tower UI for AI Agent Fleets", font_size=28, color=WHITE, bold=True)

if os.path.exists(IMG_UI):
    slide.shapes.add_picture(IMG_UI, Inches(5.8), Inches(1.6), Inches(6.8), Inches(5.2))

features = [
    ("Real-Time Policy Engine", "Intercepts & evaluates action attempts before DB write."),
    ("Cryptographic Audit Ledger", "SHA-256 parent-child hash chaining ensures tamper proofing."),
    ("Fleet Emergency Controls", "Instant per-agent or global Emergency Stop kill switch."),
    ("Spend Cap Enforcement", "Configurable daily budget limits enforced live per persona."),
]

for i, (title, desc) in enumerate(features):
    top = Inches(1.6 + i * 1.3)
    add_accent_bar(slide, Inches(0.8), top + Inches(0.05), height=Inches(0.5))
    add_text_box(slide, Inches(1.05), top, Inches(4.5), Inches(0.35),
                 title, font_size=17, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.05), top + Inches(0.38), Inches(4.5), Inches(0.8),
                 desc, font_size=13, color=GRAY)

add_slide_number(slide, 3, TOTAL_SLIDES)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SYSTEM ARCHITECTURE (with exact user white-box architecture diagram)
# ═════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_transition(slide, "fade")

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.5),
             "SYSTEM ARCHITECTURE", font_size=14, color=GOLD, bold=True)

add_text_box(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.6),
             "Inline governance between agents and execution backends",
             font_size=28, color=WHITE, bold=True)

if os.path.exists(IMG_ARCH):
    slide.shapes.add_picture(IMG_ARCH, Inches(1.2), Inches(1.6), Inches(10.9), Inches(5.3))

add_slide_number(slide, 4, TOTAL_SLIDES)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — POLICY EVALUATION FLOW
# ═════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_transition(slide, "fade")

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "EVALUATION LOGIC", font_size=14, color=GOLD, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
             "Deterministic 4-check policy pipeline", font_size=28, color=WHITE, bold=True)

steps = [
    ("1", "Agent Action Request", "Action payload sent to evaluate_policy(agent_id, action_type, amount)", BLUE_ACC),
    ("2", "Status Check", "Is agent STOPPED?  →  BLOCKED  (emergency_stop_active)", RED),
    ("3", "Block-List Check", "Is action in blocked_actions?  →  BLOCKED  (action_not_permitted)", RED),
    ("4", "Allow-List Check", "Is action missing from allowed_actions?  →  BLOCKED  (action_not_in_allowlist)", RED),
    ("5", "Spend Cap Check", "current_spend + amount > daily_cap?  →  BLOCKED  (spend_cap_exceeded)", RED),
    ("6", "Execute & Log", "Approved → Update spend → SHA-256 hash → Broadcast via WebSocket", GREEN),
]

for i, (num, title, desc, color) in enumerate(steps):
    top = Inches(1.9 + i * 0.88)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), top, Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf_c = circle.text_frame
    tf_c.paragraphs[0].text = num
    tf_c.paragraphs[0].font.size = Pt(14)
    tf_c.paragraphs[0].font.color.rgb = WHITE
    tf_c.paragraphs[0].font.bold = True
    tf_c.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_c.word_wrap = False

    add_text_box(slide, Inches(1.5), top, Inches(3.0), Inches(0.4),
                 title, font_size=17, color=WHITE, bold=True)

    add_text_box(slide, Inches(4.8), top, Inches(7.8), Inches(0.5),
                 desc, font_size=14, color=GRAY)

add_slide_number(slide, 5, TOTAL_SLIDES)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TECH STACK
# ═════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_transition(slide, "fade")

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "TECHNOLOGY STACK", font_size=14, color=GOLD, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
             "Built for low latency, auditability, and enterprise readiness",
             font_size=28, color=WHITE, bold=True)

stack_items = [
    ("Backend API", "FastAPI (Python 3.11) + Uvicorn ASGI"),
    ("Frontend UI", "React 18 + Vite + Tailwind CSS v4"),
    ("Real-Time Layer", "WebSocket (FastAPI WS Manager → React hooks)"),
    ("Crypto Audit", "SHA-256 hash chaining via Python hashlib"),
    ("Database", "SQLite (WAL mode) — swappable to PostgreSQL"),
    ("Agent Simulator", "5 persona types with configurable action profiles"),
]

for i, (label, detail) in enumerate(stack_items):
    row = i // 2
    col = i % 2
    left = Inches(0.8 + col * 6.0)
    top = Inches(2.4 + row * 1.5)

    card = add_rounded_rect(slide, left, top, Inches(5.5), Inches(1.2))

    add_text_box(slide, left + Inches(0.3), top + Inches(0.15), Inches(5.0), Inches(0.4),
                 label, font_size=18, color=GOLD, bold=True)

    add_text_box(slide, left + Inches(0.3), top + Inches(0.6), Inches(5.0), Inches(0.5),
                 detail, font_size=15, color=GRAY)

add_slide_number(slide, 6, TOTAL_SLIDES)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — BUSINESS & SOCIETAL IMPACT
# ═════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_transition(slide, "fade")

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "BUSINESS & SOCIETAL IMPACT", font_size=14, color=GOLD, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
             "Why Aegis matters for financial institutions",
             font_size=28, color=WHITE, bold=True)

impacts = [
    ("Financial Risk Mitigation",
     "Eliminates unauthorized agent spend and rogue transactions before they execute. The policy engine is the single mandatory gate across the stack.",
     "🛡"),
    ("Compliance Automation",
     "Replaces slow, manual audits with cryptographically verifiable, append-only logs aligned with PCI-DSS and SOX internal control standards.",
     "✓"),
    ("Safe AI Scaling",
     "Enables financial institutions to deploy autonomous workflows with operator controls, reducing human gating while increasing velocity responsibly.",
     "🚀"),
]

for i, (title, desc, icon) in enumerate(impacts):
    top = Inches(2.2 + i * 1.7)

    card = add_rounded_rect(slide, Inches(0.8), top, Inches(11.7), Inches(1.4))

    add_text_box(slide, Inches(1.2), top + Inches(0.15), Inches(0.6), Inches(0.5),
                 icon, font_size=24, color=GOLD)

    add_text_box(slide, Inches(1.8), top + Inches(0.15), Inches(4.0), Inches(0.4),
                 title, font_size=20, color=WHITE, bold=True)

    add_text_box(slide, Inches(1.8), top + Inches(0.65), Inches(10.2), Inches(0.7),
                 desc, font_size=14, color=GRAY)

add_slide_number(slide, 7, TOTAL_SLIDES)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — AUDIT INTEGRITY & VERIFICATION (with real Hash Audit Log screenshot)
# ═════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_transition(slide, "fade")

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.5),
             "ARCHITECTURAL GUARANTEES", font_size=14, color=GOLD, bold=True)

add_text_box(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.6),
             "Immutable SHA-256 Audit Log & Instant Chain Verification",
             font_size=28, color=WHITE, bold=True)

if os.path.exists(IMG_AUD):
    slide.shapes.add_picture(IMG_AUD, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))

add_slide_number(slide, 8, TOTAL_SLIDES)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — SCALABILITY & ASSUMPTIONS
# ═════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_transition(slide, "fade")

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "SCALABILITY & ASSUMPTIONS", font_size=14, color=GOLD, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
             "From prototype to production", font_size=28, color=WHITE, bold=True)

scale_items = [
    ("Stateless Policy Engine", "Sub-millisecond in-memory evaluation. Horizontally scalable across multiple app instances."),
    ("Database Upgrade Path", "Seamless migration from SQLite to PostgreSQL with append-only WORM storage for production."),
    ("WebSocket Pub/Sub", "Connection manager can be backed by Redis Pub/Sub for high-concurrency enterprise monitoring."),
]

for i, (title, desc) in enumerate(scale_items):
    top = Inches(2.2 + i * 1.4)
    add_accent_bar(slide, Inches(0.8), top + Inches(0.05))
    add_text_box(slide, Inches(1.15), top, Inches(5.5), Inches(0.4),
                 title, font_size=18, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.15), top + Inches(0.45), Inches(5.5), Inches(0.8),
                 desc, font_size=14, color=GRAY)

card = add_rounded_rect(slide, Inches(7.2), Inches(2.2), Inches(5.3), Inches(4.5))

add_text_box(slide, Inches(7.5), Inches(2.4), Inches(4.8), Inches(0.4),
             "Assumptions & Constraints", font_size=18, color=GOLD, bold=True)

assumptions = [
    "• Demo uses an internal multi-agent simulator via /simulate/action",
    "• Production requires agents to invoke the Aegis authorization API directly",
    "• Audit log entries are strictly append-only",
    "• Editing historical records invalidates all subsequent hashes",
    "• Emergency stop is authoritative and synchronously enforced",
]

tf = add_text_box(slide, Inches(7.5), Inches(3.0), Inches(4.8), Inches(3.5),
                  assumptions[0], font_size=13, color=GRAY)
for a in assumptions[1:]:
    add_paragraph(tf, a, font_size=13, color=GRAY, space_before=Pt(8))

add_slide_number(slide, 9, TOTAL_SLIDES)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — THANK YOU
# ═════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_transition(slide, "fade")

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = GOLD
shape.line.fill.background()

add_text_box(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1.0),
             "AEGIS", font_size=56, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.6),
             "Governance Control for Financial AI Agents",
             font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.5),
             "Thank you for your time.",
             font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = GOLD
shape.line.fill.background()

add_slide_number(slide, 10, TOTAL_SLIDES)

# ─── Save ────────────────────────────────────────────────────────────────────

prs.save(OUT_FILE)
print(f"✅ Presentation with smooth slide transitions saved to: {OUT_FILE}")
